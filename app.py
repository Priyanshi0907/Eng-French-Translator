"""
Linguo — Multilingual Translator Backend API
=============================================
Any-language translation using Meta's NLLB-200 (No Language Left Behind)
model via Hugging Face transformers. A single model handles 100+ languages
bidirectionally — no per-pair models required.

Endpoints:
  GET  /languages          - List of supported languages (code + name)
  POST /translate          - Translate text (auto-detects source if omitted)
  POST /translate-file     - Upload PDF / DOCX / PPTX, get translated text
  POST /rephrase           - AI rephrase. Three selectable engines:
                                - backtranslation (default): pivots through
                                  another language via NLLB-200. Works for
                                  any supported language, but only reshuffles
                                  what the MT model understands — struggles
                                  with typos/slang since NLLB is trained on
                                  formal text.
                                - paraphrase: a dedicated English paraphrase
                                  model (T5-based). Free/local, better at
                                  varying sentence structure, English only.
                                - llm: calls the Google Gemini API to
                                  genuinely rewrite the text, including fixing
                                  typos and slang. Needs GEMINI_API_KEY set
                                  and the `google-genai` package installed.
  GET  /health             - Health check

Run:
    pip install -r requirements.txt
    uvicorn app:app --host 0.0.0.0 --port 8000 --reload
"""

from contextlib import asynccontextmanager
from functools import lru_cache
import os
import re
import io

from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from pydantic import BaseModel, Field
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
from langdetect import detect, DetectorFactory, LangDetectException

# Make langdetect deterministic
DetectorFactory.seed = 0

# ── Language catalogue ──────────────────────────────────────────────────
# Maps a short public-facing code (used by the frontend/API) to a display
# name and the corresponding FLORES-200 code the NLLB model expects.
# Sourced from facebookresearch/flores (flores200/README.md).
LANGUAGES: dict[str, dict[str, str]] = {
    "en": {"name": "English", "flores": "eng_Latn"},
    "fr": {"name": "French", "flores": "fra_Latn"},
    "es": {"name": "Spanish", "flores": "spa_Latn"},
    "de": {"name": "German", "flores": "deu_Latn"},
    "it": {"name": "Italian", "flores": "ita_Latn"},
    "pt": {"name": "Portuguese", "flores": "por_Latn"},
    "nl": {"name": "Dutch", "flores": "nld_Latn"},
    "ru": {"name": "Russian", "flores": "rus_Cyrl"},
    "zh": {"name": "Chinese (Simplified)", "flores": "zho_Hans"},
    "zh-Hant": {"name": "Chinese (Traditional)", "flores": "zho_Hant"},
    "yue": {"name": "Cantonese", "flores": "yue_Hant"},
    "ja": {"name": "Japanese", "flores": "jpn_Jpan"},
    "ko": {"name": "Korean", "flores": "kor_Hang"},
    "ar": {"name": "Arabic", "flores": "arb_Arab"},
    "hi": {"name": "Hindi", "flores": "hin_Deva"},
    "bn": {"name": "Bengali", "flores": "ben_Beng"},
    "ur": {"name": "Urdu", "flores": "urd_Arab"},
    "fa": {"name": "Persian", "flores": "pes_Arab"},
    "tr": {"name": "Turkish", "flores": "tur_Latn"},
    "vi": {"name": "Vietnamese", "flores": "vie_Latn"},
    "th": {"name": "Thai", "flores": "tha_Thai"},
    "id": {"name": "Indonesian", "flores": "ind_Latn"},
    "ms": {"name": "Malay", "flores": "zsm_Latn"},
    "tl": {"name": "Tagalog", "flores": "tgl_Latn"},
    "sw": {"name": "Swahili", "flores": "swh_Latn"},
    "am": {"name": "Amharic", "flores": "amh_Ethi"},
    "ha": {"name": "Hausa", "flores": "hau_Latn"},
    "yo": {"name": "Yoruba", "flores": "yor_Latn"},
    "ig": {"name": "Igbo", "flores": "ibo_Latn"},
    "zu": {"name": "Zulu", "flores": "zul_Latn"},
    "xh": {"name": "Xhosa", "flores": "xho_Latn"},
    "af": {"name": "Afrikaans", "flores": "afr_Latn"},
    "sv": {"name": "Swedish", "flores": "swe_Latn"},
    "no": {"name": "Norwegian (Bokmål)", "flores": "nob_Latn"},
    "nn": {"name": "Norwegian (Nynorsk)", "flores": "nno_Latn"},
    "da": {"name": "Danish", "flores": "dan_Latn"},
    "fi": {"name": "Finnish", "flores": "fin_Latn"},
    "pl": {"name": "Polish", "flores": "pol_Latn"},
    "cs": {"name": "Czech", "flores": "ces_Latn"},
    "sk": {"name": "Slovak", "flores": "slk_Latn"},
    "hu": {"name": "Hungarian", "flores": "hun_Latn"},
    "ro": {"name": "Romanian", "flores": "ron_Latn"},
    "bg": {"name": "Bulgarian", "flores": "bul_Cyrl"},
    "uk": {"name": "Ukrainian", "flores": "ukr_Cyrl"},
    "el": {"name": "Greek", "flores": "ell_Grek"},
    "he": {"name": "Hebrew", "flores": "heb_Hebr"},
    "ka": {"name": "Georgian", "flores": "kat_Geor"},
    "hy": {"name": "Armenian", "flores": "hye_Armn"},
    "az": {"name": "Azerbaijani", "flores": "azj_Latn"},
    "kk": {"name": "Kazakh", "flores": "kaz_Cyrl"},
    "ky": {"name": "Kyrgyz", "flores": "kir_Cyrl"},
    "uz": {"name": "Uzbek", "flores": "uzn_Latn"},
    "tk": {"name": "Turkmen", "flores": "tuk_Latn"},
    "tg": {"name": "Tajik", "flores": "tgk_Cyrl"},
    "mn": {"name": "Mongolian", "flores": "khk_Cyrl"},
    "ne": {"name": "Nepali", "flores": "npi_Deva"},
    "si": {"name": "Sinhala", "flores": "sin_Sinh"},
    "ta": {"name": "Tamil", "flores": "tam_Taml"},
    "te": {"name": "Telugu", "flores": "tel_Telu"},
    "kn": {"name": "Kannada", "flores": "kan_Knda"},
    "ml": {"name": "Malayalam", "flores": "mal_Mlym"},
    "mr": {"name": "Marathi", "flores": "mar_Deva"},
    "gu": {"name": "Gujarati", "flores": "guj_Gujr"},
    "pa": {"name": "Punjabi", "flores": "pan_Guru"},
    "or": {"name": "Odia", "flores": "ory_Orya"},
    "as": {"name": "Assamese", "flores": "asm_Beng"},
    "my": {"name": "Burmese", "flores": "mya_Mymr"},
    "km": {"name": "Khmer", "flores": "khm_Khmr"},
    "lo": {"name": "Lao", "flores": "lao_Laoo"},
    "bo": {"name": "Tibetan", "flores": "bod_Tibt"},
    "dz": {"name": "Dzongkha", "flores": "dzo_Tibt"},
    "sq": {"name": "Albanian", "flores": "als_Latn"},
    "mk": {"name": "Macedonian", "flores": "mkd_Cyrl"},
    "sr": {"name": "Serbian", "flores": "srp_Cyrl"},
    "hr": {"name": "Croatian", "flores": "hrv_Latn"},
    "bs": {"name": "Bosnian", "flores": "bos_Latn"},
    "sl": {"name": "Slovenian", "flores": "slv_Latn"},
    "et": {"name": "Estonian", "flores": "est_Latn"},
    "lv": {"name": "Latvian", "flores": "lvs_Latn"},
    "lt": {"name": "Lithuanian", "flores": "lit_Latn"},
    "ga": {"name": "Irish", "flores": "gle_Latn"},
    "gd": {"name": "Scottish Gaelic", "flores": "gla_Latn"},
    "cy": {"name": "Welsh", "flores": "cym_Latn"},
    "is": {"name": "Icelandic", "flores": "isl_Latn"},
    "fo": {"name": "Faroese", "flores": "fao_Latn"},
    "mt": {"name": "Maltese", "flores": "mlt_Latn"},
    "eu": {"name": "Basque", "flores": "eus_Latn"},
    "ca": {"name": "Catalan", "flores": "cat_Latn"},
    "gl": {"name": "Galician", "flores": "glg_Latn"},
    "oc": {"name": "Occitan", "flores": "oci_Latn"},
    "eo": {"name": "Esperanto", "flores": "epo_Latn"},
    "be": {"name": "Belarusian", "flores": "bel_Cyrl"},
    "ceb": {"name": "Cebuano", "flores": "ceb_Latn"},
    "jv": {"name": "Javanese", "flores": "jav_Latn"},
    "su": {"name": "Sundanese", "flores": "sun_Latn"},
    "mg": {"name": "Malagasy", "flores": "plt_Latn"},
    "so": {"name": "Somali", "flores": "som_Latn"},
    "rw": {"name": "Kinyarwanda", "flores": "kin_Latn"},
    "rn": {"name": "Rundi", "flores": "run_Latn"},
    "ny": {"name": "Nyanja (Chichewa)", "flores": "nya_Latn"},
    "sn": {"name": "Shona", "flores": "sna_Latn"},
    "st": {"name": "Southern Sotho", "flores": "sot_Latn"},
    "tn": {"name": "Tswana", "flores": "tsn_Latn"},
    "ts": {"name": "Tsonga", "flores": "tso_Latn"},
    "ss": {"name": "Swati", "flores": "ssw_Latn"},
    "lg": {"name": "Ganda (Luganda)", "flores": "lug_Latn"},
    "ln": {"name": "Lingala", "flores": "lin_Latn"},
    "wo": {"name": "Wolof", "flores": "wol_Latn"},
    "sg": {"name": "Sango", "flores": "sag_Latn"},
    "ti": {"name": "Tigrinya", "flores": "tir_Ethi"},
    "om": {"name": "Oromo", "flores": "gaz_Latn"},
    "sd": {"name": "Sindhi", "flores": "snd_Arab"},
    "ps": {"name": "Pashto", "flores": "pbt_Arab"},
    "ku": {"name": "Kurdish (Kurmanji)", "flores": "kmr_Latn"},
    "ckb": {"name": "Kurdish (Sorani)", "flores": "ckb_Arab"},
    "yi": {"name": "Yiddish", "flores": "ydd_Hebr"},
    "sa": {"name": "Sanskrit", "flores": "san_Deva"},
    "qu": {"name": "Quechua", "flores": "quy_Latn"},
    "gn": {"name": "Guarani", "flores": "grn_Latn"},
    "ay": {"name": "Aymara", "flores": "ayr_Latn"},
    "ht": {"name": "Haitian Creole", "flores": "hat_Latn"},
    "lb": {"name": "Luxembourgish", "flores": "ltz_Latn"},
    "mi": {"name": "Maori", "flores": "mri_Latn"},
    "sm": {"name": "Samoan", "flores": "smo_Latn"},
    "fj": {"name": "Fijian", "flores": "fij_Latn"},
    "war": {"name": "Waray", "flores": "war_Latn"},
    "ilo": {"name": "Ilocano", "flores": "ilo_Latn"},
    "ast": {"name": "Asturian", "flores": "ast_Latn"},
    "scn": {"name": "Sicilian", "flores": "scn_Latn"},
}

DETECT_CODE_ALIASES = {
    # langdetect uses a couple of non-standard codes; map them onto our catalogue
    "zh-cn": "zh",
    "zh-tw": "zh-Hant",
}

DEFAULT_SOURCE = "en"
DEFAULT_TARGET = "fr"

MODEL_NAME = "facebook/nllb-200-distilled-600M"
MAX_INPUT_LENGTH = 500        # chars for /translate
MAX_FILE_TEXT    = 15_000     # chars extracted from uploaded files

# ── Model loading ──────────────────────────────────────────────────────
tokenizer = None
model = None


def load_model():
    global tokenizer, model
    tokenizer = AutoTokenizer.from_pretrained(MODEL_NAME)
    model = AutoModelForSeq2SeqLM.from_pretrained(MODEL_NAME)
    model.eval()


# ── Rephrase engine #2: dedicated English paraphrase model ─────────────
# Lazy-loaded only if someone actually picks the "paraphrase" engine, so a
# default deployment doesn't pay for a second model download/load.
PARAPHRASE_MODEL_NAME = "humarin/chatgpt_paraphraser_on_T5_base"
_paraphrase_tokenizer = None
_paraphrase_model = None


def _load_paraphrase_model():
    global _paraphrase_tokenizer, _paraphrase_model
    if _paraphrase_model is None:
        _paraphrase_tokenizer = AutoTokenizer.from_pretrained(PARAPHRASE_MODEL_NAME)
        _paraphrase_model = AutoModelForSeq2SeqLM.from_pretrained(PARAPHRASE_MODEL_NAME)
        _paraphrase_model.eval()


def paraphrase_with_t5(text: str) -> str:
    _load_paraphrase_model()
    input_ids = _paraphrase_tokenizer(
        f"paraphrase: {text}",
        return_tensors="pt",
        padding="longest",
        max_length=128,
        truncation=True,
    ).input_ids
    outputs = _paraphrase_model.generate(
        input_ids,
        temperature=0.7,
        repetition_penalty=10.0,
        num_return_sequences=1,
        no_repeat_ngram_size=2,
        num_beams=5,
        num_beam_groups=5,
        max_length=128,
        diversity_penalty=3.0,
    )
    return _paraphrase_tokenizer.batch_decode(outputs, skip_special_tokens=True)[0]


# ── Rephrase engine #3: LLM (Google Gemini API) ─────────────────────────
# Genuinely understands intent, typos and slang — unlike the two MT-based
# engines above, which can only rearrange words the translation model
# already recognizes. Requires GEMINI_API_KEY and the `google-genai` package.
_gemini_client = None
_gemini_import_error = None
try:
    from google import genai as google_genai
    if os.environ.get("GEMINI_API_KEY"):
        _gemini_client = google_genai.Client(api_key=os.environ["GEMINI_API_KEY"])
except ImportError as e:
    _gemini_import_error = str(e)

LLM_MODEL_NAME = "gemini-2.5-flash"


def rephrase_with_llm(text: str, lang_name: str) -> str:
    if _gemini_client is None:
        reason = (
            "the 'google-genai' package is not installed"
            if _gemini_import_error
            else "GEMINI_API_KEY is not set"
        )
        raise HTTPException(
            status_code=503,
            detail=f"LLM rephrase engine is not configured ({reason}).",
        )
    prompt = (
        f"Rephrase the following {lang_name} text naturally, fixing any typos, "
        f"abbreviations, or grammar issues, while preserving its meaning and tone. "
        f"Reply with ONLY the rephrased text, nothing else.\n\nText: {text}"
    )
    response = _gemini_client.models.generate_content(
        model=LLM_MODEL_NAME,
        contents=prompt,
    )
    return response.text.strip()


@asynccontextmanager
async def lifespan(app: FastAPI):
    print(f"Loading {MODEL_NAME}... (first run downloads ~2.4GB)")
    load_model()
    print(f"Model loaded. Ready to translate {len(LANGUAGES)} languages.")
    yield


app = FastAPI(title="Linguo Multilingual Translator API", lifespan=lifespan)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.get("/")
def serve_index():
    return FileResponse("index.html")


# ── Schemas ────────────────────────────────────────────────────────────
class TranslateRequest(BaseModel):
    text: str = Field(..., min_length=1, max_length=MAX_INPUT_LENGTH)
    source_lang: str | None = Field(
        default=None,
        description="Language code (see /languages). Omit to auto-detect.",
    )
    target_lang: str | None = Field(
        default=None,
        description="Language code (see /languages). Defaults to French, or "
                    "English if the source is French.",
    )


class TranslateResponse(BaseModel):
    translated_text: str
    source_lang: str
    target_lang: str


class RephraseRequest(BaseModel):
    text: str = Field(..., min_length=1, max_length=MAX_INPUT_LENGTH)
    source_lang: str | None = Field(
        default=None,
        description="Language code (see /languages). Omit to auto-detect.",
    )
    engine: str = Field(
        default="backtranslation",
        description="One of: 'backtranslation' (any language), "
                    "'paraphrase' (English only, T5-based), "
                    "'llm' (any language, needs GEMINI_API_KEY).",
    )


class RephraseResponse(BaseModel):
    original: str
    rephrased: str
    engine: str


class FileTranslateResponse(BaseModel):
    filename: str
    source_lang: str
    target_lang: str
    translated_text: str
    char_count: int


class LanguageInfo(BaseModel):
    code: str
    name: str


# ── Language helpers ───────────────────────────────────────────────────
def validate_lang(code: str) -> str:
    if code not in LANGUAGES:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported language code '{code}'. See /languages for the full list.",
        )
    return code


def detect_language(text: str) -> str:
    try:
        code = detect(text)
    except LangDetectException:
        return DEFAULT_SOURCE
    code = DETECT_CODE_ALIASES.get(code, code)
    return code if code in LANGUAGES else DEFAULT_SOURCE


def resolve_target(source_lang: str, target_lang: str | None) -> str:
    if target_lang is not None:
        return validate_lang(target_lang)
    # Backward-compatible default: opposite of EN/FR, else English.
    if source_lang == "en":
        return "fr"
    return "en"


# ── Core translation helpers ───────────────────────────────────────────
@lru_cache(maxsize=1024)
def translate_cached(text: str, source_lang: str, target_lang: str) -> str:
    src_flores = LANGUAGES[source_lang]["flores"]
    tgt_flores = LANGUAGES[target_lang]["flores"]

    tokenizer.src_lang = src_flores
    batch = tokenizer(text, return_tensors="pt", padding=True, truncation=True, max_length=512)
    forced_bos_token_id = tokenizer.convert_tokens_to_ids(tgt_flores)
    generated = model.generate(
        **batch,
        forced_bos_token_id=forced_bos_token_id,
        max_new_tokens=512,
    )
    return tokenizer.batch_decode(generated, skip_special_tokens=True)[0]


def split_sentences(text: str) -> list[str]:
    parts = re.split(r'(?<=[.!?])\s+', text.strip())
    return [p.strip() for p in parts if p.strip()]


def translate_full(text: str, source_lang: str, target_lang: str) -> str:
    """Translate by splitting into sentences so nothing is silently dropped."""
    sentences = split_sentences(text)
    if len(sentences) <= 1:
        return translate_cached(text, source_lang, target_lang)
    return " ".join(translate_cached(s, source_lang, target_lang) for s in sentences)


def translate_preserving_layout(text: str, source_lang: str, target_lang: str) -> str:
    """Translates paragraph-by-paragraph and line-by-line to preserve layout,
    line breaks, and document flow."""
    paragraphs = text.split("\n\n")
    translated_paragraphs = []
    for para in paragraphs:
        if not para.strip():
            translated_paragraphs.append("")
            continue
        lines = para.split("\n")
        translated_lines = []
        for line in lines:
            if not line.strip():
                translated_lines.append("")
                continue
            translated_lines.append(translate_full(line, source_lang, target_lang))
        translated_paragraphs.append("\n".join(translated_lines))
    return "\n\n".join(translated_paragraphs)


# ── File text extraction ───────────────────────────────────────────────
def extract_pdf(data: bytes) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            pages = [page.extract_text() or "" for page in pdf.pages]
        return "\n\n".join(pages).strip()
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Could not read PDF: {e}")


def extract_docx(data: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs).strip()
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Could not read DOCX: {e}")


def extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        slides = []
        for slide in prs.slides:
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if texts:
                slides.append("\n".join(texts))
        return "\n\n".join(slides).strip()
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Could not read PPTX: {e}")


EXTRACTORS = {
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
}


# ── Endpoints ──────────────────────────────────────────────────────────
@app.get("/languages", response_model=list[LanguageInfo])
def list_languages():
    return sorted(
        [{"code": code, "name": info["name"]} for code, info in LANGUAGES.items()],
        key=lambda x: x["name"],
    )


@app.post("/translate", response_model=TranslateResponse)
def translate(request: TranslateRequest):
    text = request.text.strip()
    if not text:
        raise HTTPException(status_code=400, detail="Text must not be empty.")

    source_lang = request.source_lang
    if source_lang is None:
        source_lang = detect_language(text)
    else:
        source_lang = validate_lang(source_lang)

    target_lang = resolve_target(source_lang, request.target_lang)

    if source_lang == target_lang:
        raise HTTPException(status_code=400, detail="Source and target languages must differ.")

    try:
        translated = translate_full(text, source_lang, target_lang)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Translation failed: {e}")

    return TranslateResponse(
        translated_text=translated,
        source_lang=source_lang,
        target_lang=target_lang,
    )


@app.post("/translate-file", response_model=FileTranslateResponse)
async def translate_file(
    file: UploadFile = File(...),
    source_lang: str | None = Form(default=None),
    target_lang: str | None = Form(default=None),
):
    """Accept a PDF, DOCX or PPTX file and return the translated text."""
    filename = file.filename or ""
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext not in EXTRACTORS:
        raise HTTPException(
            status_code=415,
            detail=f"Unsupported file type '{ext}'. Supported: .pdf, .docx, .pptx",
        )

    data = await file.read()
    raw_text = EXTRACTORS[ext](data)

    if not raw_text:
        raise HTTPException(status_code=422, detail="No text found in the file.")

    # Truncate to avoid huge latency
    if len(raw_text) > MAX_FILE_TEXT:
        raw_text = raw_text[:MAX_FILE_TEXT] + "\n\n[Truncated — file too large]"

    # Detect / validate language
    if source_lang is None:
        source_lang = detect_language(raw_text[:500])
    else:
        source_lang = validate_lang(source_lang)

    target_lang = resolve_target(source_lang, target_lang)

    if source_lang == target_lang:
        raise HTTPException(status_code=400, detail="Source and target languages must differ.")

    try:
        translated = translate_preserving_layout(raw_text, source_lang, target_lang)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Translation failed: {e}")

    return FileTranslateResponse(
        filename=filename,
        source_lang=source_lang,
        target_lang=target_lang,
        translated_text=translated,
        char_count=len(raw_text),
    )


class DocumentGenerateRequest(BaseModel):
    text: str
    filename: str


@app.post("/generate-docx")
def generate_docx(request: DocumentGenerateRequest):
    try:
        from docx import Document
        doc = Document()

        # Style the docx neatly
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Arial'
        font.size = 11

        paragraphs = request.text.split("\n\n")
        for para in paragraphs:
            if para.strip():
                p = doc.add_paragraph()
                lines = para.split("\n")
                for idx, line in enumerate(lines):
                    if idx > 0:
                        p.add_run("\n")
                    p.add_run(line)
            else:
                doc.add_paragraph()

        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)

        headers = {
            "Content-Disposition": f"attachment; filename={request.filename}",
            "Access-Control-Expose-Headers": "Content-Disposition"
        }
        return StreamingResponse(
            buffer,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers=headers
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate Word doc: {e}")


@app.post("/generate-pdf")
def generate_pdf(request: DocumentGenerateRequest):
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_LEFT

        buffer = io.BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=54,
            leftMargin=54,
            topMargin=54,
            bottomMargin=54
        )

        styles = getSampleStyleSheet()
        body_style = ParagraphStyle(
            'ElegantBody',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=10.5,
            leading=15,
            textColor='#453F57',
            alignment=TA_LEFT,
            spaceAfter=10
        )

        story = []
        paragraphs = request.text.split("\n\n")
        for para in paragraphs:
            if para.strip():
                formatted_text = para.replace("\n", "<br/>")
                p = Paragraph(formatted_text, body_style)
                story.append(p)
            else:
                story.append(Spacer(1, 10))

        doc.build(story)
        buffer.seek(0)

        headers = {
            "Content-Disposition": f"attachment; filename={request.filename}",
            "Access-Control-Expose-Headers": "Content-Disposition"
        }
        return StreamingResponse(
            buffer,
            media_type="application/pdf",
            headers=headers
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate PDF: {e}")


@app.post("/rephrase", response_model=RephraseResponse)
def rephrase(request: RephraseRequest):
    """Rephrase text using one of three engines:
    - backtranslation: EN<->pivot MT round-trip (any language, weak on typos/slang)
    - paraphrase: dedicated T5 paraphrase model (English only)
    - llm: Claude via the Anthropic API (any language, handles typos/slang)
    """
    text = request.text.strip()
    if not text:
        raise HTTPException(status_code=400, detail="Text must not be empty.")

    source_lang = request.source_lang or detect_language(text)
    source_lang = validate_lang(source_lang)

    engine = (request.engine or "backtranslation").lower()

    try:
        if engine == "llm":
            rephrased = rephrase_with_llm(text, LANGUAGES[source_lang]["name"])
        elif engine == "paraphrase":
            if source_lang != "en":
                raise HTTPException(
                    status_code=400,
                    detail="The paraphrase engine only supports English. "
                           "Use 'backtranslation' or 'llm' for other languages.",
                )
            rephrased = paraphrase_with_t5(text)
        elif engine == "backtranslation":
            pivot_lang = "fr" if source_lang == "en" else "en"
            pivot     = translate_full(text, source_lang, pivot_lang)   # e.g. EN -> FR
            rephrased = translate_full(pivot, pivot_lang, source_lang)  # e.g. FR -> EN
        else:
            raise HTTPException(status_code=400, detail=f"Unknown engine '{engine}'.")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Rephrase failed: {e}")

    return RephraseResponse(original=text, rephrased=rephrased, engine=engine)


@app.get("/health")
def health():
    return {
        "status": "ok",
        "model_loaded": model is not None,
        "languages_supported": len(LANGUAGES),
        "llm_engine_available": _gemini_client is not None,
    }